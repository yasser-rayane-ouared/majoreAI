"""
StudyFlow AI - File Processor Service
Extracts text from various file types for RAG indexing.
"""
import os
import io
import csv
import zipfile
from pathlib import Path
from typing import Dict, List, Optional


async def extract_text_from_file(file_path: str, original_name: str, mime_type: str = "") -> Dict:
    """Extract text content from a file based on its type."""
    ext = Path(original_name).suffix.lower()
    result = {
        "file_name": original_name,
        "file_type": "unknown",
        "text": "",
        "error": None,
    }

    try:
        # PDF
        if ext == ".pdf":
            result["file_type"] = "PDF Document"
            result["text"] = await _extract_pdf(file_path)

        # DOCX
        elif ext in (".docx", ".doc"):
            result["file_type"] = "Word Document"
            result["text"] = await _extract_docx(file_path)

        # PPTX
        elif ext in (".pptx", ".ppt"):
            result["file_type"] = "PowerPoint"
            result["text"] = await _extract_pptx(file_path)

        # CSV
        elif ext == ".csv":
            result["file_type"] = "CSV Spreadsheet"
            result["text"] = await _extract_csv(file_path)

        # Images (OCR)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
            result["file_type"] = "Image (OCR)"
            result["text"] = await _extract_image_ocr(file_path)

        # ZIP
        elif ext in (".zip",):
            result["file_type"] = "ZIP Archive"
            result["text"] = await _extract_zip(file_path, original_name)

        # Source code files
        elif ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c", ".cs",
                      ".go", ".rb", ".php", ".swift", ".kt", ".rs", ".sql", ".r",
                      ".html", ".css", ".json", ".xml", ".yaml", ".yml", ".sh", ".bash"):
            lang = ext.lstrip(".")
            result["file_type"] = f"Code ({lang.upper()})"
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            result["text"] = f"```{lang}\n{content}\n```"

        # Plain text / Markdown
        elif ext in (".txt", ".md", ".log", ".text"):
            result["file_type"] = "Text Document"
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                result["text"] = f.read()

        # Fallback: try reading as text
        else:
            result["file_type"] = f"File ({ext})"
            try:
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    result["text"] = f.read()
            except Exception:
                result["text"] = f"[Could not extract text from {original_name}]"

    except Exception as e:
        result["error"] = str(e)
        result["text"] = f"[Error processing {original_name}: {str(e)}]"

    return result


async def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyPDF2."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        return "\n\n".join(pages) if pages else "[PDF contained no extractable text]"
    except ImportError:
        return "[PyPDF2 not installed - run: pip install PyPDF2]"
    except Exception as e:
        return f"[PDF extraction error: {str(e)}]"


async def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        
        return "\n\n".join(paragraphs) if paragraphs else "[DOCX contained no text]"
    except ImportError:
        return "[python-docx not installed - run: pip install python-docx]"
    except Exception as e:
        return f"[DOCX extraction error: {str(e)}]"


async def _extract_pptx(file_path: str) -> str:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
        return "\n\n".join(slides) if slides else "[PPTX contained no text]"
    except ImportError:
        return "[python-pptx not installed - run: pip install python-pptx]"
    except Exception as e:
        return f"[PPTX extraction error: {str(e)}]"


async def _extract_csv(file_path: str) -> str:
    """Extract text from CSV."""
    try:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i > 500:  # Limit rows
                    rows.append("... (truncated)")
                    break
                rows.append(" | ".join(row))
        return "\n".join(rows) if rows else "[CSV is empty]"
    except Exception as e:
        return f"[CSV extraction error: {str(e)}]"


async def _extract_image_ocr(file_path: str) -> str:
    """Extract text from image using OCR."""
    try:
        import pytesseract
        from PIL import Image
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image, lang="eng+fra+ara")
        return text.strip() if text.strip() else "[No text detected in image]"
    except ImportError:
        # Fallback: just describe the image
        return "[Image uploaded - OCR requires pytesseract. Install with: pip install pytesseract Pillow]"
    except Exception as e:
        return f"[OCR error: {str(e)}]"


async def _extract_zip(file_path: str, original_name: str) -> str:
    """Extract and process files within a ZIP archive."""
    try:
        extracted_texts = []
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in zf.namelist()[:50]:  # Limit to 50 files
                if name.endswith("/"):
                    continue
                ext = Path(name).suffix.lower()
                try:
                    with zf.open(name) as f:
                        content = f.read()
                        try:
                            text = content.decode("utf-8", errors="replace")
                            extracted_texts.append(f"--- {name} ---\n{text[:5000]}")
                        except Exception:
                            extracted_texts.append(f"--- {name} --- [binary file]")
                except Exception as e:
                    extracted_texts.append(f"--- {name} --- [error: {str(e)}]")
        
        return "\n\n".join(extracted_texts) if extracted_texts else "[ZIP is empty]"
    except Exception as e:
        return f"[ZIP extraction error: {str(e)}]"


def get_file_icon(file_type: str) -> str:
    """Get an emoji icon for a file type."""
    icons = {
        "PDF Document": "📄",
        "Word Document": "📝",
        "PowerPoint": "📊",
        "CSV Spreadsheet": "📈",
        "Image (OCR)": "🖼️",
        "ZIP Archive": "📦",
        "Text Document": "📃",
    }
    if "Code" in file_type:
        return "💻"
    return icons.get(file_type, "📎")
